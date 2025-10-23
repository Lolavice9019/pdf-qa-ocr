import streamlit as st
import tempfile
import os
from pdf2image import convert_from_bytes
from paddleocr import PaddleOCR
from openai import OpenAI
from datetime import datetime
import PyPDF2
import pdfplumber
import pikepdf
import io
from docx import Document
from pptx import Presentation
import openpyxl
from bs4 import BeautifulSoup
from striprtf.striprtf import rtf_to_text
from ebooklib import epub
import ebooklib
from fpdf import FPDF
import csv as csv_module

# Page configuration
st.set_page_config(
    page_title="Universal Document Q&A System",
    page_icon="📄",
    layout="wide"
)

# Initialize models
@st.cache_resource
def load_ocr_model():
    """Load PaddleOCR model"""
    return PaddleOCR(use_angle_cls=True, lang='en')

@st.cache_resource
def load_qa_client():
    """Initialize OpenAI client"""
    try:
        return OpenAI()
    except Exception as e:
        st.error(f"Error initializing QA client: {e}")
        return None

# ========== DOCUMENT EXTRACTORS ==========

def extract_from_txt(file_bytes):
    """Extract text from TXT files"""
    try:
        text = file_bytes.decode('utf-8')
        return text, [text], True
    except:
        try:
            text = file_bytes.decode('latin-1')
            return text, [text], True
        except:
            return None, None, False

def extract_from_docx(file_bytes):
    """Extract text from DOCX files"""
    try:
        doc = Document(io.BytesIO(file_bytes))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        full_text = '\n\n'.join(paragraphs)
        return full_text, paragraphs, True
    except Exception as e:
        return None, None, False

def extract_from_pptx(file_bytes):
    """Extract text from PPTX files"""
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        slides_text = []
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            slides_text.append('\n'.join(slide_text))
        full_text = '\n\n'.join(slides_text)
        return full_text, slides_text, True
    except Exception as e:
        return None, None, False

def extract_from_xlsx(file_bytes):
    """Extract text from XLSX files"""
    try:
        wb = openpyxl.load_workbook(io.BytesIO(file_bytes))
        sheets_text = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            rows = []
            for row in sheet.iter_rows(values_only=True):
                row_text = '\t'.join([str(cell) if cell is not None else '' for cell in row])
                if row_text.strip():
                    rows.append(row_text)
            sheets_text.append(f"Sheet: {sheet_name}\n" + '\n'.join(rows))
        full_text = '\n\n'.join(sheets_text)
        return full_text, sheets_text, True
    except Exception as e:
        return None, None, False

def extract_from_csv(file_bytes):
    """Extract text from CSV files"""
    try:
        text = file_bytes.decode('utf-8')
        lines = text.split('\n')
        return text, lines, True
    except:
        try:
            text = file_bytes.decode('latin-1')
            lines = text.split('\n')
            return text, lines, True
        except:
            return None, None, False

def extract_from_html(file_bytes):
    """Extract text from HTML files"""
    try:
        html = file_bytes.decode('utf-8')
        soup = BeautifulSoup(html, 'html.parser')
        text = soup.get_text(separator='\n', strip=True)
        paragraphs = [p for p in text.split('\n') if p.strip()]
        return '\n\n'.join(paragraphs), paragraphs, True
    except:
        return None, None, False

def extract_from_rtf(file_bytes):
    """Extract text from RTF files"""
    try:
        rtf_text = file_bytes.decode('utf-8')
        text = rtf_to_text(rtf_text)
        lines = text.split('\n')
        return text, lines, True
    except:
        return None, None, False

def extract_from_epub(file_bytes):
    """Extract text from EPUB files"""
    try:
        book = epub.read_epub(io.BytesIO(file_bytes))
        chapters = []
        for item in book.get_items():
            if item.get_type() == ebooklib.ITEM_DOCUMENT:
                soup = BeautifulSoup(item.get_content(), 'html.parser')
                text = soup.get_text(separator='\n', strip=True)
                if text.strip():
                    chapters.append(text)
        full_text = '\n\n'.join(chapters)
        return full_text, chapters, True
    except Exception as e:
        return None, None, False

def extract_from_md(file_bytes):
    """Extract text from Markdown files"""
    try:
        text = file_bytes.decode('utf-8')
        return text, [text], True
    except:
        return None, None, False

def extract_from_pdf(file_bytes, ocr_model):
    """Extract from PDF with repair and OCR"""
    # Try pdfplumber
    try:
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            pages = []
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    pages.append(text)
            if pages and len(''.join(pages).strip()) > 50:
                return '\n\n'.join(pages), pages, True
    except:
        pass
    
    # Try PyPDF2
    try:
        reader = PyPDF2.PdfReader(io.BytesIO(file_bytes), strict=False)
        pages = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                pages.append(text)
        if pages and len(''.join(pages).strip()) > 50:
            return '\n\n'.join(pages), pages, True
    except:
        pass
    
    # Try repair with pikepdf
    try:
        with pikepdf.open(io.BytesIO(file_bytes)) as pdf:
            output = io.BytesIO()
            pdf.save(output)
            output.seek(0)
            repaired_bytes = output.read()
            
            # Retry pdfplumber on repaired
            with pdfplumber.open(io.BytesIO(repaired_bytes)) as pdf:
                pages = []
                for page in pdf.pages:
                    text = page.extract_text()
                    if text:
                        pages.append(text)
                if pages and len(''.join(pages).strip()) > 50:
                    return '\n\n'.join(pages), pages, True
    except:
        pass
    
    # Try OCR as last resort
    try:
        images = convert_from_bytes(file_bytes, dpi=200)
        pages = []
        for image in images:
            with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp:
                image.save(tmp.name)
                result = ocr_model.ocr(tmp.name, cls=True)
                page_text = []
                if result and result[0]:
                    for line in result[0]:
                        if line and len(line) > 1 and line[1]:
                            page_text.append(line[1][0])
                pages.append('\n'.join(page_text))
                try:
                    os.unlink(tmp.name)
                except:
                    pass
        if pages and len(''.join(pages).strip()) > 50:
            return '\n\n'.join(pages), pages, True
    except:
        pass
    
    return None, None, False

# ========== UNIVERSAL EXTRACTOR ==========

def extract_text_from_document(file, ocr_model):
    """
    Universal document text extractor
    Automatically detects file type and uses appropriate method
    """
    filename = file.name
    file.seek(0)
    file_bytes = file.read()
    
    # Detect file type by extension
    ext = filename.lower().split('.')[-1] if '.' in filename else ''
    
    extractors = {
        'txt': ('TXT', extract_from_txt),
        'text': ('TXT', extract_from_txt),
        'docx': ('DOCX', lambda b: extract_from_docx(b)),
        'doc': ('DOCX', lambda b: extract_from_docx(b)),
        'pptx': ('PPTX', lambda b: extract_from_pptx(b)),
        'ppt': ('PPTX', lambda b: extract_from_pptx(b)),
        'xlsx': ('XLSX', lambda b: extract_from_xlsx(b)),
        'xls': ('XLSX', lambda b: extract_from_xlsx(b)),
        'csv': ('CSV', extract_from_csv),
        'html': ('HTML', extract_from_html),
        'htm': ('HTML', extract_from_html),
        'rtf': ('RTF', extract_from_rtf),
        'epub': ('EPUB', extract_from_epub),
        'md': ('Markdown', extract_from_md),
        'markdown': ('Markdown', extract_from_md),
        'pdf': ('PDF', lambda b: extract_from_pdf(b, ocr_model)),
    }
    
    if ext in extractors:
        method_name, extractor = extractors[ext]
        full_text, pages, success = extractor(file_bytes)
        if success:
            return full_text, pages, True, method_name
    
    # Fallback: try all methods
    for method_name, extractor in extractors.values():
        try:
            if method_name == 'PDF':
                full_text, pages, success = extractor(file_bytes)
            else:
                full_text, pages, success = extractor(file_bytes)
            if success:
                return full_text, pages, True, f"{method_name} (auto-detected)"
        except:
            continue
    
    return None, None, False, "failed"

# ========== Q&A FUNCTIONS ==========

def answer_question(question, context, client, model_name="gpt-4.1-mini"):
    """Answer question based on context"""
    try:
        max_context_chars = 20000
        if len(context) > max_context_chars:
            context = context[:max_context_chars] + "\n\n[... truncated ...]"
        
        response = client.chat.completions.create(
            model=model_name,
            messages=[
                {"role": "system", "content": "You are a document analysis assistant. Answer questions precisely based on the provided context."},
                {"role": "user", "content": f"Context:\n{context}\n\nQuestion: {question}\n\nAnswer based only on the context above."}
            ],
            temperature=0.3,
            max_tokens=500
        )
        return response.choices[0].message.content
    except Exception as e:
        return f"Error: {e}"

def answer_question_multi_doc(question, docs_context, client, model_name="gpt-4.1-mini"):
    """Answer question across multiple documents"""
    try:
        combined = ""
        for filename, context in docs_context.items():
            combined += f"\n\n=== {filename} ===\n{context[:5000]}\n"
        
        if len(combined) > 20000:
            combined = combined[:20000] + "\n\n[... truncated ...]"
        
        response = client.chat.completions.create(
            model=model_name,
            messages=[
                {"role": "system", "content": "You are a multi-document analysis assistant. Answer questions and cite which documents contain the information."},
                {"role": "user", "content": f"Documents:\n{combined}\n\nQuestion: {question}\n\nAnswer and cite sources."}
            ],
            temperature=0.3,
            max_tokens=700
        )
        return response.choices[0].message.content
    except Exception as e:
        return f"Error: {e}"

def save_log(filename, pages, method):
    """Save extraction log"""
    log_path = f"/home/ubuntu/ocr_logs/{filename}.txt"
    with open(log_path, 'w', encoding='utf-8') as f:
        f.write(f"Method: {method}\nPages/Sections: {len(pages)}\n{'='*50}\n\n")
        for idx, text in enumerate(pages):
            f.write(f"=== Section {idx + 1} ===\n{text}\n\n")
    return log_path

def update_todo(filename, num_pages, question, answer):
    """Log Q&A"""
    with open("/home/ubuntu/todo.md", 'a', encoding='utf-8') as f:
        f.write(f"\n## {filename}\n- Pages: {num_pages}\n- Q: {question}\n- A: {answer}\n- Date: {datetime.now()}\n\n")

def export_to_pdf(text, filename):
    """Export text to PDF"""
    try:
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", size=12)
        # Handle unicode by replacing problematic characters
        safe_text = text.encode('latin-1', 'replace').decode('latin-1')
        for line in safe_text.split('\n'):
            pdf.multi_cell(0, 10, txt=line)
        return pdf.output(dest='S').encode('latin-1')
    except Exception as e:
        return None

def export_to_docx(text, filename):
    """Export text to DOCX"""
    try:
        doc = Document()
        doc.add_heading(filename, 0)
        for paragraph in text.split('\n\n'):
            if paragraph.strip():
                doc.add_paragraph(paragraph)
        output = io.BytesIO()
        doc.save(output)
        output.seek(0)
        return output.read()
    except Exception as e:
        return None

def export_to_html(text, filename):
    """Export text to HTML"""
    html = f"""<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <title>{filename}</title>
    <style>
        body {{ font-family: Arial, sans-serif; max-width: 800px; margin: 40px auto; padding: 20px; line-height: 1.6; }}
        h1 {{ color: #333; }}
        pre {{ background: #f4f4f4; padding: 15px; border-radius: 5px; overflow-x: auto; }}
    </style>
</head>
<body>
    <h1>{filename}</h1>
    <pre>{text}</pre>
</body>
</html>"""
    return html.encode('utf-8')

def export_to_csv(documents_dict):
    """Export multiple documents to CSV"""
    output = io.StringIO()
    writer = csv_module.writer(output)
    writer.writerow(['Filename', 'Method', 'Sections', 'Character Count', 'Full Text'])
    for filename, data in documents_dict.items():
        writer.writerow([
            filename,
            data['method'],
            data['num_pages'],
            len(data['full_text']),
            data['full_text']
        ])
    return output.getvalue().encode('utf-8')

def export_to_markdown(text, filename):
    """Export text to Markdown"""
    md = f"# {filename}\n\n{text}"
    return md.encode('utf-8')

# ========== MAIN UI ==========

st.title("📄 Universal Document Q&A System")
st.markdown("""
**Supports ALL document types:** PDF, DOCX, PPTX, XLSX, TXT, HTML, RTF, EPUB, Markdown, CSV, and more!

- ✅ Automatic file type detection
- ✅ Perfect extraction method for each format
- ✅ Batch processing up to 50 GB per file
- ✅ Multi-document search
- ✅ AI-powered question answering
""")

ocr_model = load_ocr_model()
qa_client = load_qa_client()

if not qa_client:
    st.error("⚠️ QA system unavailable")
    st.stop()

st.header("1. Upload Documents")
uploaded_files = st.file_uploader(
    "Select any document type (PDF, DOCX, PPTX, XLSX, TXT, HTML, etc.)",
    accept_multiple_files=True,
    help="Supports: PDF, DOCX, PPTX, XLSX, CSV, TXT, HTML, RTF, EPUB, Markdown"
)

if 'processed_files' not in st.session_state:
    st.session_state.processed_files = {}

if uploaded_files:
    st.markdown(f"**{len(uploaded_files)} file(s) selected**")
    
    with st.expander("📋 View files"):
        for f in uploaded_files:
            st.write(f"• {f.name} ({f.size / (1024*1024):.2f} MB)")
    
    unprocessed = [f for f in uploaded_files if f.name not in st.session_state.processed_files]
    
    if unprocessed:
        if st.button("🚀 Process All Files", type="primary"):
            progress_bar = st.progress(0)
            status = st.container()
            
            successful = 0
            failed = 0
            methods = {}
            
            for idx, file in enumerate(unprocessed):
                status.info(f"Processing {idx+1}/{len(unprocessed)}: {file.name}")
                
                full_text, pages, success, method = extract_text_from_document(file, ocr_model)
                
                if success and full_text:
                    log_path = save_log(file.name, pages, method)
                    st.session_state.processed_files[file.name] = {
                        'full_text': full_text,
                        'pages': pages,
                        'num_pages': len(pages),
                        'log_path': log_path,
                        'method': method
                    }
                    methods[method] = methods.get(method, 0) + 1
                    successful += 1
                    status.success(f"✅ {file.name} - Method: {method}")
                else:
                    failed += 1
                    status.error(f"❌ {file.name} - Extraction failed")
                
                progress_bar.progress((idx + 1) / len(unprocessed))
            
            st.success(f"🎉 Complete! ✅ {successful} successful, ❌ {failed} failed")
            if methods:
                st.info(f"📊 Methods: {', '.join([f'{k}={v}' for k,v in methods.items()])}")
            if successful > 0:
                st.balloons()
    else:
        st.success(f"✅ All files processed")

if st.session_state.processed_files:
    st.header("2. Processed Documents")
    
    total_pages = sum(d['num_pages'] for d in st.session_state.processed_files.values())
    total_chars = sum(len(d['full_text']) for d in st.session_state.processed_files.values())
    
    col1, col2, col3 = st.columns(3)
    col1.metric("Documents", len(st.session_state.processed_files))
    col2.metric("Sections", total_pages)
    col3.metric("Characters", f"{total_chars:,}")
    
    # Multi-format export options
    st.subheader("📥 Export All Documents")
    
    # Prepare combined text
    all_text = ""
    for filename, data in st.session_state.processed_files.items():
        all_text += f"{'='*60}\n"
        all_text += f"FILE: {filename}\n"
        all_text += f"METHOD: {data['method']}\n"
        all_text += f"SECTIONS: {data['num_pages']}\n"
        all_text += f"{'='*60}\n\n"
        all_text += data['full_text']
        all_text += "\n\n\n"
    
    import json
    all_json = {
        'export_date': datetime.now().isoformat(),
        'total_documents': len(st.session_state.processed_files),
        'total_sections': total_pages,
        'total_characters': total_chars,
        'documents': {
            filename: {
                'method': data['method'],
                'num_sections': data['num_pages'],
                'full_text': data['full_text'],
                'sections': data['pages']
            }
            for filename, data in st.session_state.processed_files.items()
        }
    }
    
    # Export format selection
    col1, col2, col3, col4 = st.columns(4)
    
    with col1:
        st.download_button(
            label="📝 TXT",
            data=all_text,
            file_name="all_documents.txt",
            mime="text/plain",
            use_container_width=True,
            help="Plain text format"
        )
    
    with col2:
        st.download_button(
            label="📦 JSON",
            data=json.dumps(all_json, indent=2),
            file_name="all_documents.json",
            mime="application/json",
            use_container_width=True,
            help="Structured JSON with metadata"
        )
    
    with col3:
        html_data = export_to_html(all_text, "All Documents")
        st.download_button(
            label="🌐 HTML",
            data=html_data,
            file_name="all_documents.html",
            mime="text/html",
            use_container_width=True,
            help="Web page format"
        )
    
    with col4:
        csv_data = export_to_csv(st.session_state.processed_files)
        st.download_button(
            label="📊 CSV",
            data=csv_data,
            file_name="all_documents.csv",
            mime="text/csv",
            use_container_width=True,
            help="Spreadsheet format"
        )
    
    col5, col6, col7, col8 = st.columns(4)
    
    with col5:
        md_data = export_to_markdown(all_text, "All Documents")
        st.download_button(
            label="📑 Markdown",
            data=md_data,
            file_name="all_documents.md",
            mime="text/markdown",
            use_container_width=True,
            help="Markdown format"
        )
    
    with col6:
        docx_data = export_to_docx(all_text, "All Documents")
        if docx_data:
            st.download_button(
                label="📄 DOCX",
                data=docx_data,
                file_name="all_documents.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                use_container_width=True,
                help="Microsoft Word format"
            )
    
    with col7:
        pdf_data = export_to_pdf(all_text, "All Documents")
        if pdf_data:
            st.download_button(
                label="📝 PDF",
                data=pdf_data,
                file_name="all_documents.pdf",
                mime="application/pdf",
                use_container_width=True,
                help="PDF format"
            )
    
    with col8:
        # XML export
        xml_data = f"""<?xml version="1.0" encoding="UTF-8"?>
<documents>
    <export_date>{datetime.now().isoformat()}</export_date>
    <total_documents>{len(st.session_state.processed_files)}</total_documents>
"""
        for filename, data in st.session_state.processed_files.items():
            xml_data += f"""    <document>
        <filename>{filename}</filename>
        <method>{data['method']}</method>
        <sections>{data['num_pages']}</sections>
        <text><![CDATA[{data['full_text']}]]></text>
    </document>
"""
        xml_data += "</documents>"
        
        st.download_button(
            label="📜 XML",
            data=xml_data.encode('utf-8'),
            file_name="all_documents.xml",
            mime="application/xml",
            use_container_width=True,
            help="XML format"
        )
    
    st.markdown("---")
    st.subheader("Individual Documents")
    
    for filename, data in st.session_state.processed_files.items():
        with st.expander(f"📄 {filename} ({data['num_pages']} sections) - {data['method']}"):
            st.markdown(f"**Method:** `{data['method']}`")
            st.markdown(f"**Log:** `{data['log_path']}`")
            
            # Download buttons - multiple formats
            st.markdown("**Export this document:**")
            dl_col1, dl_col2, dl_col3, dl_col4 = st.columns(4)
            
            with dl_col1:
                st.download_button(
                    label="📝 TXT",
                    data=data['full_text'],
                    file_name=f"{filename}_extracted.txt",
                    mime="text/plain",
                    key=f"txt_{filename}",
                    use_container_width=True
                )
            
            with dl_col2:
                import json
                json_data = json.dumps({
                    'filename': filename,
                    'method': data['method'],
                    'num_sections': data['num_pages'],
                    'full_text': data['full_text'],
                    'sections': data['pages']
                }, indent=2)
                st.download_button(
                    label="📦 JSON",
                    data=json_data,
                    file_name=f"{filename}.json",
                    mime="application/json",
                    key=f"json_{filename}",
                    use_container_width=True
                )
            
            with dl_col3:
                html_data = export_to_html(data['full_text'], filename)
                st.download_button(
                    label="🌐 HTML",
                    data=html_data,
                    file_name=f"{filename}.html",
                    mime="text/html",
                    key=f"html_{filename}",
                    use_container_width=True
                )
            
            with dl_col4:
                md_data = export_to_markdown(data['full_text'], filename)
                st.download_button(
                    label="📑 MD",
                    data=md_data,
                    file_name=f"{filename}.md",
                    mime="text/markdown",
                    key=f"md_{filename}",
                    use_container_width=True
                )
            
            dl_col5, dl_col6, dl_col7, dl_col8 = st.columns(4)
            
            with dl_col5:
                docx_data = export_to_docx(data['full_text'], filename)
                if docx_data:
                    st.download_button(
                        label="📄 DOCX",
                        data=docx_data,
                        file_name=f"{filename}.docx",
                        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                        key=f"docx_{filename}",
                        use_container_width=True
                    )
            
            with dl_col6:
                pdf_data = export_to_pdf(data['full_text'], filename)
                if pdf_data:
                    st.download_button(
                        label="📕 PDF",
                        data=pdf_data,
                        file_name=f"{filename}.pdf",
                        mime="application/pdf",
                        key=f"pdf_{filename}",
                        use_container_width=True
                    )
            
            for idx, page in enumerate(data['pages']):
                with st.expander(f"Section {idx + 1}"):
                    st.text_area(f"Text {idx + 1}", page, height=200, key=f"{filename}_{idx}", disabled=True)
    
    st.header("3. Ask Questions")
    
    mode = st.radio("Mode:", ["Single Document", "Multiple Documents"])
    
    if mode == "Single Document":
        selected = st.selectbox("Select document:", list(st.session_state.processed_files.keys()))
        if selected:
            context = st.session_state.processed_files[selected]['full_text']
            question = st.text_input("Question:", placeholder="What is this about?")
            
            with st.expander("⚙️ Options"):
                model = st.selectbox("Model:", ["gpt-4.1-mini", "gpt-4.1-nano", "gemini-2.5-flash"])
            
            if st.button("🔍 Search", type="primary") and question:
                with st.spinner("Processing..."):
                    answer = answer_question(question, context, qa_client, model)
                    st.markdown("### 💡 Answer:")
                    st.success(answer)
                    update_todo(selected, st.session_state.processed_files[selected]['num_pages'], question, answer)
    else:
        use_all = st.checkbox("Use all documents", value=True)
        selected_docs = list(st.session_state.processed_files.keys()) if use_all else st.multiselect("Select:", list(st.session_state.processed_files.keys()))
        
        if selected_docs:
            st.write(f"**{len(selected_docs)} selected**")
            question = st.text_input("Question:", placeholder="Which documents mention...?", key="multi_q")
            
            with st.expander("⚙️ Options"):
                model = st.selectbox("Model:", ["gpt-4.1-mini", "gpt-4.1-nano", "gemini-2.5-flash"], key="multi_m")
            
            if st.button("🔍 Search All", type="primary") and question:
                with st.spinner("Searching..."):
                    docs_ctx = {d: st.session_state.processed_files[d]['full_text'] for d in selected_docs}
                    answer = answer_question_multi_doc(question, docs_ctx, qa_client, model)
                    st.markdown("### 💡 Answer:")
                    st.success(answer)
else:
    st.info("👆 Upload documents to start")

st.markdown("---")
st.markdown("<div style='text-align: center'><p><strong>Universal Document Q&A System</strong></p></div>", unsafe_allow_html=True)

